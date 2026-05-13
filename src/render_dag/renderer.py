"""
renderer.py
---------------
PowerPoint rendering library for Airflow DAG diagrams.

The "content" (what is connected to what) lives in the YAML; the
"appearance" (colors, fonts, dimensions) lives in theme.py; the
"rendering logic" (layout calculation, line drawing) lives in this file.
A three-layer structure.

Main features:
  - Auto layout: divides the slide width by the max column count, auto-tunes
    box size and font size, and centers the used row range vertically inside
    the content area.
  - Status management: three badge types + styles (new / changed / deleted).
  - Elbow (right-angled) connectors: the arrowhead is drawn only on the
    terminal segment.
  - Strips page-number placeholders (removes unexpected placeholders that
    sometimes remain on the blank layout so they don't auto-appear when
    PowerPoint opens the file).
"""

from __future__ import annotations
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree

from . import theme as t


# ---- Layout ----------------------------------------------------------------
class Layout:
    """Grid dimensions + font scale."""
    def __init__(self, origin_x=0.5, origin_y=1.0,
                 col_w=2.3, row_h=1.0, box_w=2.0, box_h=0.75,
                 font_scale=1.0):
        self.origin_x   = origin_x
        self.origin_y   = origin_y
        self.col_w      = col_w
        self.row_h      = row_h
        self.box_w      = box_w
        self.box_h      = box_h
        self.font_scale = font_scale

    def cell_to_xy(self, col, row):
        return (self.origin_x + col * self.col_w,
                self.origin_y + row * self.row_h)

def compute_auto_layout(slide_def):
    """Compute an auto layout from the slide's boxes.
    Strategy:
      - Column width: divide the slide width evenly by the max column count.
      - Rows: fixed row height; center the used row range vertically inside
        the content area.
      - Font: scale proportionally with box_w / DEFAULT_BOX_W
        (clamped between min/max).
    """
    boxes = slide_def.get("boxes", [])
    grid_boxes = [b for b in boxes if "col" in b and "row" in b]
    if not grid_boxes:
        return {}

    max_col    = max(b["col"] + b.get("colspan", 1) - 1 for b in grid_boxes)
    cols       = max_col + 1
    rows_used  = sorted({b["row"] for b in grid_boxes})
    min_row    = rows_used[0]
    max_row    = rows_used[-1]
    n_rows     = max_row - min_row + 1

    # Horizontal: divide the slide width evenly by the max column count.
    available_w = t.SLIDE_W_IN - 2 * t.AUTO_MARGIN_X
    col_w       = available_w / cols
    box_w       = col_w * (1 - t.AUTO_COL_GAP_RATIO)

    font_scale = max(t.AUTO_FONT_SCALE_MIN,
                     min(t.AUTO_FONT_SCALE_MAX, box_w / t.DEFAULT_BOX_W))

    # Vertical: center the rows inside the content area.
    has_footer = bool(slide_def.get("header"))
    if has_footer:
        # The lead sits just under the title underline, so the DAG starts below it.
        content_top = t.FOOTER_Y + t.FOOTER_H + t.CONTENT_PADDING
    else:
        content_top = t.TITLE_BAR_BOTTOM + t.CONTENT_PADDING
    content_bottom = t.SLIDE_H_IN - t.CONTENT_PADDING
    available_h    = content_bottom - content_top

    box_h = max(0.55, min(0.85, 0.62 + 0.15 * (font_scale - 1)))

    # When there are many rows, shrink box_h to keep a minimum gap so boxes
    # don't overlap.
    if n_rows > 1:
        MIN_ROW_GAP        = 0.18
        max_box_h_with_gap = (available_h - (n_rows - 1) * MIN_ROW_GAP) / n_rows
        if max_box_h_with_gap < box_h:
            box_h = max(0.40, max_box_h_with_gap)

    ideal_row_h = box_h * t.AUTO_ROW_H_MULTIPLIER
    if n_rows > 1:
        max_row_h = (available_h - box_h) / (n_rows - 1)
        row_h     = min(ideal_row_h, max_row_h)
    else:
        row_h = ideal_row_h

    used_h   = (n_rows - 1) * row_h + box_h
    top_y    = content_top + (available_h - used_h) / 2
    origin_y = top_y - min_row * row_h

    return {
        "origin_x":   t.AUTO_MARGIN_X,
        "origin_y":   origin_y,
        "col_w":      col_w,
        "row_h":      row_h,
        "box_w":      box_w,
        "box_h":      box_h,
        "font_scale": font_scale,
    }

# ---- Helpers ---------------------------------------------------------------
def inch(v):
    return Emu(int(v * 914400))


def _set_font(run, size=10, bold=False, color=None, font=None):
    if font is None:
        font = t.FONT
    run.font.name      = font
    run.font.size      = Pt(float(size))
    run.font.bold      = bold
    if color is not None:
        run.font.color.rgb = color
    rPr = run._r.get_or_add_rPr()
    for tag in ("eastAsia", "cs"):
        a_tag = f"a:{tag}"
        el = rPr.find(qn(a_tag))
        if el is None:
            el = etree.SubElement(rPr, qn(a_tag))
        el.set("typeface", font)


def _set_dashed_line(shape):
    ln       = shape.line._get_or_add_ln()
    existing = ln.find(qn("a:prstDash"))
    if existing is not None:
        ln.remove(existing)
    prstDash = etree.SubElement(ln, qn("a:prstDash"))
    prstDash.set("val", "dash")


def _strip_layout_placeholders(prs):
    """Remove the date / footer / sldNum placeholders from the blank layout
    and master. Prevents PowerPoint from auto-inserting them on open."""
    target_types = {"dt", "ftr", "sldNum"}
    targets = list(prs.slide_layouts) + list(prs.slide_masters)
    for container in targets:
        spTree = container._element.find(
            qn("p:cSld") + "/" + qn("p:spTree")
        )
        if spTree is None:
            continue
        for sp in list(spTree.findall(qn("p:sp"))):
            nvSpPr = sp.find(qn("p:nvSpPr"))
            if nvSpPr is None:
                continue
            nvPr = nvSpPr.find(qn("p:nvPr"))
            if nvPr is None:
                continue
            ph = nvPr.find(qn("p:ph"))
            if ph is None:
                continue
            if ph.get("type") in target_types:
                spTree.remove(sp)


# ---- Drawing primitives ----------------------------------------------------
def add_title_bar(slide, title):
    """Left-aligned title + underline (uniform design across slides)."""
    textbox = slide.shapes.add_textbox(
        inch(t.TITLE_X), inch(t.TITLE_Y),
        inch(t.SLIDE_W_IN - 2 * t.TITLE_X), inch(t.TITLE_H)
    )
    tf                    = textbox.text_frame
    tf.margin_left        = Emu(0)
    tf.margin_right       = Emu(0)
    tf.margin_top         = Emu(0)
    tf.margin_bottom      = Emu(0)
    tf.word_wrap          = True
    tf.vertical_anchor    = MSO_ANCHOR.MIDDLE
    p                     = tf.paragraphs[0]
    p.alignment           = PP_ALIGN.LEFT
    p.text                = ""
    run                   = p.add_run()
    run.text              = title
    _set_font(run, size=t.TITLE_FONT_SIZE, bold=True, color=t.TITLE_TEXT_COLOR)

    # Underline.
    line = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        inch(t.TITLE_X), inch(t.TITLE_LINE_Y),
        inch(t.SLIDE_W_IN - t.TITLE_X), inch(t.TITLE_LINE_Y)
    )
    line.line.color.rgb = t.TITLE_LINE_COLOR
    line.line.width     = Pt(t.TITLE_LINE_WEIGHT)


def add_header(slide, text):
    # Lead text placed right under the title underline. No frame, no background, text only.
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        inch(t.FOOTER_X), inch(t.FOOTER_Y),
        inch(t.SLIDE_W_IN - 2 * t.FOOTER_X), inch(t.FOOTER_H)
    )
    shape.fill.background()
    shape.line.fill.background()
    # Remove all shadow / effect filters.
    spPr = shape._element.find(qn("p:spPr"))
    if spPr is not None:
        for tag in (qn("a:effectLst"), qn("a:effectDag")):
            el = spPr.find(tag)
            if el is not None:
                spPr.remove(el)
        effectLst = etree.SubElement(spPr, qn("a:effectLst"))
    tf                   = shape.text_frame
    tf.word_wrap         = True
    tf.margin_left       = tf.margin_right  = Emu(0)
    tf.margin_top        = tf.margin_bottom = Emu(0)
    tf.text              = ""
    for i, line in enumerate(text.split("\n")):
        p           = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        run         = p.add_run()
        run.text    = line
        _set_font(run, size=t.FOOTER_FONT_SIZE, color=t.FOOTER_TEXT)


def add_task_box(slide, layout, b):
    if "x_in" in b and "y_in" in b:
        x, y = b["x_in"], b["y_in"]
    else:
        x, y = layout.cell_to_xy(b["col"], b["row"])

    colspan = b.get("colspan", 1)
    w       = b.get("w_in", layout.box_w + (colspan - 1) * layout.col_w)
    h       = b.get("h_in", layout.box_h)
    status  = b.get("status")

    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        inch(x), inch(y), inch(w), inch(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = (
        t.BOX_DELETED_FILL if status == "deleted" else t.BOX_FILL
    )
    shape.line.color.rgb  = t.BOX_LINE
    shape.line.width      = Pt(t.BOX_LINE_WEIGHT)
    if status == "deleted":
        _set_dashed_line(shape)

    tf                    = shape.text_frame
    tf.margin_left        = tf.margin_right  = Emu(45720)
    tf.margin_top         = tf.margin_bottom = Emu(18288)
    tf.word_wrap          = True
    tf.vertical_anchor    = MSO_ANCHOR.MIDDLE
    tf.text               = ""

    label_size    = b.get("label_size",    t.BOX_LABEL_SIZE)    * layout.font_scale
    sublabel_size = b.get("sublabel_size", t.BOX_SUBLABEL_SIZE) * layout.font_scale

    p           = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run         = p.add_run()
    run.text    = b["label"]
    _set_font(run, size=label_size, bold=True, color=t.BOX_TEXT)

    if b.get("sublabel"):
        p2          = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        run2         = p2.add_run()
        run2.text    = b["sublabel"]
        _set_font(run2, size=sublabel_size, color=t.BOX_TEXT)

    badge = None
    # Status badge.
    if status in t.BADGE_COLORS:
        badge_text = t.BADGE_TEXTS[status]
        bw         = 0.65 if badge_text == "CHANGED" else 0.55
        bx         = x + w - bw + 0.05
        by         = y - 0.14
        badge      = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            inch(bx), inch(by), inch(bw), inch(t.BADGE_H)
        )
        badge.fill.solid()
        badge.fill.fore_color.rgb = t.BADGE_COLORS[status]
        badge.line.fill.background()
        btf                       = badge.text_frame
        btf.margin_left           = btf.margin_right  = Emu(0)
        btf.margin_top            = btf.margin_bottom = Emu(0)
        btf.vertical_anchor       = MSO_ANCHOR.MIDDLE
        bp                        = btf.paragraphs[0]
        bp.alignment              = PP_ALIGN.CENTER
        brun                      = bp.add_run()
        brun.text                 = badge_text
        _set_font(brun, size=t.BADGE_FONT_SIZE, bold=True,
                        color=t.BADGE_TEXT_COLOR)

    return {"x": x, "y": y, "w": w, "h": h, "shape": shape, "badge": badge}


def add_label(slide, layout, lbl):
    """Free-position text. Used for callouts, headings, and annotations.

    Coordinates:
      - (col, row): grid-relative (moves with the boxes; row may be fractional)
      - (x_in, y_in): absolute inches from the top-left of the slide
    """
    if "col" in lbl and "row" in lbl:
        cx = layout.origin_x + lbl["col"] * layout.col_w + layout.box_w / 2
        cy = layout.origin_y + lbl["row"] * layout.row_h + layout.box_h / 2
        w  = lbl.get("w_in", layout.box_w * 0.9)
        h  = lbl.get("h_in", 0.6)
        x  = cx - w / 2
        y  = cy - h / 2
    else:
        x = lbl["x_in"]
        y = lbl["y_in"]
        w = lbl.get("w_in", 2.0)
        h = lbl.get("h_in", 0.4)

    bg     = lbl.get("bg_color")
    border = lbl.get("border_color")
    style  = lbl.get("style")
    if style == "callout":
        bg     = bg     or t.LABEL_CALLOUT_BG
        border = border or t.LABEL_CALLOUT_LINE

    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        inch(x), inch(y), inch(w), inch(h)
    )
    if bg:
        shape.fill.solid()
        shape.fill.fore_color.rgb = bg
    else:
        shape.fill.background()
    if border:
        shape.line.color.rgb = border
        shape.line.width     = Pt(lbl.get("border_weight", 1.0))
        if lbl.get("border_dash"):
            _set_dashed_line(shape)
    else:
        shape.line.fill.background()

    tf                    = shape.text_frame
    tf.margin_left        = tf.margin_right  = Emu(45720)
    tf.margin_top         = tf.margin_bottom = Emu(22860)
    tf.word_wrap          = True
    tf.vertical_anchor    = MSO_ANCHOR.MIDDLE

    align_map = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER,
                 "right": PP_ALIGN.RIGHT}
    alignment = align_map.get(lbl.get("align", "center"), PP_ALIGN.CENTER)

    tf.text = ""
    for i, line in enumerate(lbl["text"].split("\n")):
        p           = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = alignment
        run         = p.add_run()
        run.text    = line
        _set_font(
            run,
            size=lbl.get("size", t.LABEL_DEFAULT_SIZE),
            bold=lbl.get("bold", False),
            color=lbl.get("text_color", RGBColor(0x20, 0x20, 0x20)),
        )


# ==================== Native connector + hidden anchors ====================
# An anchor is a tiny shape with a black-line border. We attach connection info
# via direct XML, so moving an anchor in PowerPoint drags the arrow with it
# (the dot mark indicates a connected endpoint).
#
# Connection-site indices (rectangle), measured against python-pptx:
#   idx 0 = top-center, idx 1 = left-center, idx 2 = bottom-center, idx 3 = right-center
#
# Anchor placement strategy:
#   For sources: pinned just inside the right edge of the box (anchor's right
#                edge = box's right edge) -> the arrow exits from site 3 (right).
#   For targets: pinned just inside the left edge of the box (anchor's left
#                edge = box's left edge) -> the arrow tip arrives at site 1
#                (left) and visually pierces the left edge of the box.
#
# At the end we move the box and badge to the end of spTree, so the box is
# brought to the front and the anchor (black border) is fully hidden behind it.
# We do NOT call begin_connect / end_connect (they perturb the connector's
# coordinates). Instead we add <a:stCxn> / <a:endCxn> as raw XML to set only
# the connection info.
_ANCHOR_SIZE = 0.10


def _add_source_anchor(slide, box_right, y):
    """Source anchor: pinned just inside the right edge of the box. The anchor's
    right edge coincides with the box's right edge. Since the box ends up in
    front, the anchor is fully hidden. The connector exits from the right edge
    (site 3)."""
    size = _ANCHOR_SIZE
    a = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        inch(box_right - size), inch(y - size / 2),
        inch(size), inch(size))
    a.fill.background()
    a.line.color.rgb = RGBColor(0x00, 0x00, 0x00)
    a.line.width     = Pt(0.5)
    return a


def _add_target_anchor(slide, box_left, y):
    """Target anchor: pinned just inside the left edge of the box. The anchor's
    left edge coincides with the box's left edge. Since the box ends up in
    front, the anchor is fully hidden. The connector tip (>) arrives at the
    left edge (site 1) and appears to pierce the left edge of the box."""
    size = _ANCHOR_SIZE
    a = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        inch(box_left), inch(y - size / 2),
        inch(size), inch(size))
    a.fill.background()
    a.line.color.rgb = RGBColor(0x00, 0x00, 0x00)
    a.line.width     = Pt(0.5)
    return a


def _attach_connector_xml(conn, src_anchor, dst_anchor):
    """Attach connection info via direct XML, without calling begin_connect /
    end_connect. begin_connect would move the connector's coordinates, which
    we want to avoid.
    idx=3 is right-center (source exit); idx=1 is left-center (target entry)."""
    nvCxnSpPr  = conn._element.find(qn("p:nvCxnSpPr"))
    cNvCxnSpPr = nvCxnSpPr.find(qn("p:cNvCxnSpPr"))
    for tag in (qn("a:stCxn"), qn("a:endCxn")):
        el = cNvCxnSpPr.find(tag)
        if el is not None:
            cNvCxnSpPr.remove(el)
    st = etree.SubElement(cNvCxnSpPr, qn("a:stCxn"))
    st.set("id",  str(src_anchor.shape_id))
    st.set("idx", "3")   # right-center
    en = etree.SubElement(cNvCxnSpPr, qn("a:endCxn"))
    en.set("id",  str(dst_anchor.shape_id))
    en.set("idx", "1")   # left-center


def _add_elbow_connector(slide, src_anchor, dst_anchor,
                         with_arrow=True, adj1=None):
    """Connect two anchor shapes with a native elbow connector.

    - Exits from src_anchor's right edge (site 3) and enters dst_anchor's left
      edge (site 1).
    - Since each anchor is hidden at the inner edge of its box, the arrow
      appears to leave the right edge of the source box and arrive at the
      left edge of the destination box.
    - The arrowhead (tailEnd) is drawn at the connector's end (target side).
    """
    # Source side: anchor's right edge = box's right edge.
    src_x = src_anchor.left + src_anchor.width
    src_y = src_anchor.top  + src_anchor.height // 2
    # Target side: anchor's left edge = box's left edge.
    dst_x = dst_anchor.left
    dst_y = dst_anchor.top  + dst_anchor.height // 2

    conn = slide.shapes.add_connector(
        MSO_CONNECTOR.ELBOW,
        src_x, src_y, dst_x, dst_y)

    # Set the connection info via direct XML (begin_connect / end_connect
    # would break the connector's position).
    _attach_connector_xml(conn, src_anchor, dst_anchor)

    conn.line.color.rgb = t.CONNECTOR_COLOR
    conn.line.width     = Pt(t.CONNECTOR_WEIGHT)
    if with_arrow:
        ln      = conn.line._get_or_add_ln()
        tailEnd = ln.find(qn("a:tailEnd"))
        if tailEnd is None:
            tailEnd = etree.SubElement(ln, qn("a:tailEnd"))
        tailEnd.set("type", "triangle")
        tailEnd.set("w", "lg")
        tailEnd.set("h", "lg")
        headEnd = ln.find(qn("a:headEnd"))
        if headEnd is not None:
            ln.remove(headEnd)

    pstyle = conn._element.find(qn("p:style"))
    if pstyle is not None:
        conn._element.remove(pstyle)

    if adj1 is not None:
        spPr     = conn._element.find(qn("p:spPr"))
        prstGeom = spPr.find(qn("a:prstGeom"))
        avLst    = prstGeom.find(qn("a:avLst"))
        if avLst is None:
            avLst = etree.SubElement(prstGeom, qn("a:avLst"))
        for gd in list(avLst):
            avLst.remove(gd)
        gd = etree.SubElement(avLst, qn("a:gd"))
        gd.set("name", "adj1")
        gd.set("fmla", f"val {int(adj1 * 100000)}")

    return conn


def compute_box_position(layout, b):
    """Compute the position of a box from its definition (used before drawing)."""
    if "x_in" in b and "y_in" in b:
        x, y = b["x_in"], b["y_in"]
    else:
        x, y = layout.cell_to_xy(b["col"], b["row"])
    colspan = b.get("colspan", 1)
    w       = b.get("w_in", layout.box_w + (colspan - 1) * layout.col_w)
    h       = b.get("h_in", layout.box_h)
    return {"x": x, "y": y, "w": w, "h": h}


# ---- High-level API --------------------------------------------------------
def render_slide(prs, slide_def):
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)

    add_title_bar(slide, slide_def["title"])

    auto = compute_auto_layout(slide_def)

    def pick(key, default):
        if key in slide_def:
            return slide_def[key]
        if key in auto:
            return auto[key]
        return default

    layout = Layout(
        origin_x=pick("origin_x",   0.5),
        origin_y=pick("origin_y",   1.0),
        col_w=pick("col_w",         2.3),
        row_h=pick("row_h",         1.0),
        box_w=pick("box_w",         2.0),
        box_h=pick("box_h",         0.75),
        font_scale=pick("font_scale", 1.0),
    )

    # ---- Pre-compute box positions (needed for endpoint splitting) ----
    positions = {}
    for b in slide_def.get("boxes", []):
        positions[b["id"]] = compute_box_position(layout, b)

    # ---- Draw boxes ----
    box_refs = {}
    for b in slide_def.get("boxes", []):
        box_refs[b["id"]] = add_task_box(slide, layout, b)

    # ---- Arrow endpoint splitting: when a box has multiple outgoing/incoming
    # edges, spread the endpoint y-coordinates. ----
    edges = slide_def.get("edges", [])
    from collections import defaultdict

    out_edges = defaultdict(list)  # src_id -> [(edge_idx, dst_id), ...]
    in_edges  = defaultdict(list)  # dst_id -> [(edge_idx, src_id), ...]
    for idx, (s, d) in enumerate(edges):
        out_edges[s].append((idx, d))
        in_edges[d].append((idx, s))

    exit_y, entry_y = {}, {}

    # Fan-out side: sort by ascending dst y, then assign exit_y top to bottom.
    #   -> Arrows aimed at upper rows leave the top of the box; arrows aimed
    #      at lower rows leave the bottom.
    for src_id, outs in out_edges.items():
        src = box_refs[src_id]
        n   = len(outs)
        if n == 1 or not t.ENDPOINT_SPLIT_ENABLED:
            for idx, _ in outs:
                exit_y[idx] = src["y"] + src["h"] / 2
        else:
            outs_sorted = sorted(outs,
                key=lambda e: box_refs[e[1]]["y"] + box_refs[e[1]]["h"] / 2)
            for i, (idx, _) in enumerate(outs_sorted):
                exit_y[idx] = src["y"] + src["h"] * (i + 1) / (n + 1)

    # Fan-in side: sort by ascending src y, then assign entry_y top to bottom.
    for dst_id, ins in in_edges.items():
        dst = box_refs[dst_id]
        m   = len(ins)
        if m == 1 or not t.ENDPOINT_SPLIT_ENABLED:
            for idx, _ in ins:
                entry_y[idx] = dst["y"] + dst["h"] / 2
        else:
            ins_sorted = sorted(ins,
                key=lambda e: box_refs[e[1]]["y"] + box_refs[e[1]]["h"] / 2)
            for j, (idx, _) in enumerate(ins_sorted):
                entry_y[idx] = dst["y"] + dst["h"] * (j + 1) / (m + 1)

    # ---- Vertical pairing: compute adj1 (the position of the middle vertical
    # segment, 0-1) per connector. ----
    # In fan-out / fan-in, upward and downward arrows don't overlap in y, so
    # we align them on the same vertical lane for symmetry.
    # adj1 = the in-connector x position of the vertical segment (0 = start,
    #        1 = end).
    # Paired edges sharing the same gap have equal connector bounding boxes,
    # so equal adj1 = equal absolute x position.
    gap_groups = defaultdict(list)
    for idx, (s, d) in enumerate(edges):
        src, dst = box_refs[s], box_refs[d]
        y1, y2   = exit_y[idx], entry_y[idx]
        if abs(y1 - y2) < 0.01:
            continue   # Same row -> straight line; no adj1 needed.
        gap_key = (round(src["x"] + src["w"], 2), round(dst["x"], 2))
        gap_groups[gap_key].append((idx, y1, y2))

    edge_adj1 = {}

    def _place_paired(upper_ids, lower_ids):
        n_ranks = max(len(upper_ids), len(lower_ids))
        if n_ranks == 0:
            return
        for rank in range(n_ranks):
            adj1 = (rank + 1) / (n_ranks + 1)
            if rank < len(upper_ids):
                edge_adj1[upper_ids[rank]] = adj1
            if rank < len(lower_ids):
                edge_adj1[lower_ids[rank]] = adj1

    for gap_key, items in gap_groups.items():
        gap_edge_ids    = [idx for idx, _, _ in items]
        sources_in_gap  = {edges[idx][0] for idx in gap_edge_ids}
        targets_in_gap  = {edges[idx][1] for idx in gap_edge_ids}
        n               = len(items)

        if len(sources_in_gap) == 1 and n >= 2:
            # Pure fan-out: sort by distance (farther first), pair same ranks.
            upper = [(idx, abs(y2 - y1)) for idx, y1, y2 in items if y2 < y1]
            lower = [(idx, abs(y2 - y1)) for idx, y1, y2 in items if y2 > y1]
            upper.sort(key=lambda e: -e[1])
            lower.sort(key=lambda e: -e[1])
            _place_paired([idx for idx, _ in upper],
                          [idx for idx, _ in lower])
        elif len(targets_in_gap) == 1 and n >= 2:
            # Pure fan-in: sort by distance (closer first), pair same ranks.
            upper = [(idx, abs(y2 - y1)) for idx, y1, y2 in items if y1 < y2]
            lower = [(idx, abs(y2 - y1)) for idx, y1, y2 in items if y1 > y2]
            upper.sort(key=lambda e: e[1])
            lower.sort(key=lambda e: e[1])
            _place_paired([idx for idx, _ in upper],
                          [idx for idx, _ in lower])
        else:
            # Mixed: spread lanes per source.
            source_groups = defaultdict(list)
            for idx in gap_edge_ids:
                source_groups[edges[idx][0]].append(idx)
            sources_sorted = sorted(source_groups.keys(),
                                    key=lambda s: box_refs[s]["y"])
            ordered = []
            for src_id in sources_sorted:
                ordered.extend(source_groups[src_id])
            for lane, idx in enumerate(ordered):
                edge_adj1[idx] = (lane + 1) / (n + 1)

    # ---- Draw edges (native elbow connector + visible anchors + adj1 control) ----
    # Place a source anchor pinned to the outside of the source's right edge,
    # and a target anchor pinned at the destination's left edge. Connect them
    # with a native elbow connector.
    # adj1 controls the x position of the middle vertical segment, lining up
    # upper / lower pairs on the same vertical lane.
    for idx, (s, d) in enumerate(edges):
        src, dst = box_refs[s], box_refs[d]
        y1       = exit_y[idx]
        y2       = entry_y[idx]
        src_anchor = _add_source_anchor(slide, src["x"] + src["w"], y1)
        dst_anchor = _add_target_anchor(slide, dst["x"], y2)

        adj1 = edge_adj1.get(idx, None)
        _add_elbow_connector(slide, src_anchor, dst_anchor,
                             with_arrow=True, adj1=adj1)

    for lbl in slide_def.get("labels", []):
        add_label(slide, layout, lbl)

    if slide_def.get("header"):
        add_header(slide, slide_def["header"])

    # ---- z-order tweak: bring task boxes (and their badges) to the front. ----
    # This hides the anchors (black borders) behind the boxes, so each arrow
    # endpoint appears to attach exactly to the edge of the box.
    spTree = slide.shapes._spTree
    for ref in box_refs.values():
        main = ref.get("shape")
        if main is not None:
            elem   = main._element
            parent = elem.getparent()
            if parent is not None:
                parent.remove(elem)
                parent.append(elem)
        badge = ref.get("badge")
        if badge is not None:
            elem   = badge._element
            parent = elem.getparent()
            if parent is not None:
                parent.remove(elem)
                parent.append(elem)


def render_deck(deck, path):
    prs              = Presentation()
    prs.slide_width  = inch(t.SLIDE_W_IN)
    prs.slide_height = inch(t.SLIDE_H_IN)

    # Strip page-number / date / footer placeholders.
    _strip_layout_placeholders(prs)

    for s in deck:
        render_slide(prs, s)
    prs.save(path)


# ---- YAML loading ----------------------------------------------------------
def _parse_color(v):
    if v is None:
        return None
    if isinstance(v, RGBColor):
        return v
    if isinstance(v, str):
        if v in t.NAMED_COLORS:
            return t.NAMED_COLORS[v]
        s = v.lstrip("#")
        if len(s) == 6:
            return RGBColor(int(s[0:2], 16), int(s[2:4], 16), int(s[4:6], 16))
    raise ValueError(f"Unrecognized color: {v!r}")


def _parse_stage(stage, expr):
    s = stage.strip()
    if not s:
        raise ValueError(f"Empty stage in edge expression: {expr!r}")
    if s.startswith("[") and s.endswith("]"):
        inner = s[1:-1]
        ids   = [x.strip() for x in inner.split(",")]
        if not all(ids) or not ids:
            raise ValueError(f"Empty ID in list stage: {expr!r}")
        return ids
    if "[" in s or "]" in s or "," in s:
        raise ValueError(f"Malformed stage {s!r} in: {expr!r}")
    return [s]


def _parse_edge_expression(expr):
    """Expand an Airflow-style edge expression into a sequence of (src, dst) tuples.

    Examples:
        "a >> b"           -> [(a, b)]
        "a >> b >> c"      -> [(a, b), (b, c)]
        "a >> [b, c]"      -> [(a, b), (a, c)]
        "[a, b] >> c"      -> [(a, c), (b, c)]
        "a >> [b, c] >> d" -> [(a, b), (a, c), (b, d), (c, d)]
    """
    stages = [s for s in expr.split(">>")]
    if len(stages) < 2:
        raise ValueError(f"Edge expression must contain '>>': {expr!r}")
    parsed = [_parse_stage(stage, expr) for stage in stages]
    edges  = []
    for i in range(len(parsed) - 1):
        for src in parsed[i]:
            for dst in parsed[i + 1]:
                edges.append((src, dst))
    return edges


def _normalize_slide(slide):
    for lbl in slide.get("labels", []) or []:
        for key in ("text_color", "bg_color", "border_color"):
            if key in lbl:
                lbl[key] = _parse_color(lbl[key])
    if "edges" in slide:
        normalized = []
        for e in slide["edges"]:
            if isinstance(e, str):
                normalized.extend(_parse_edge_expression(e))
            else:
                normalized.append(tuple(e))
        slide["edges"] = normalized
    return slide


def load_deck_from_yaml(path):
    import yaml
    with open(path, encoding="utf-8") as f:
        data = yaml.safe_load(f)
    if not isinstance(data, dict) or "slides" not in data:
        raise ValueError(f"{path}: top-level must have a 'slides' key")
    return [_normalize_slide(s) for s in data["slides"]]
